"""Bounded extraction and OCR for supported Drive artifacts."""

from __future__ import annotations

import csv
from datetime import date, datetime, time
import json
import math
from pathlib import Path, PurePosixPath
import re
import stat
import struct
import time as clock
from typing import Mapping
import unicodedata
import warnings
from zipfile import BadZipFile, ZipFile

from bs4 import BeautifulSoup
from docx import Document
import fitz
from openpyxl import load_workbook
from openpyxl.utils.exceptions import InvalidFileException
from PIL import Image, ImageOps, UnidentifiedImageError
from pptx import Presentation
import pytesseract

from .models import ExtractedBlock, ExtractedDocument
from .protocol import DriveRagError


MAX_FILE_BYTES = 250 * 1024 * 1024
MAX_SHEET_CELLS = 50_000
MAX_JSON_NODES = 20_000
MAX_STRUCTURED_ITEMS = 20_000
OCR_TIMEOUT_SECONDS = 600
PDF_OCR_TEXT_THRESHOLD = 40
PDF_OCR_DPI = 200
MAX_PDF_PAGES = 10_000
MAX_PDF_FILE_BYTES = 32 * 1024 * 1024
MAX_PDF_OCR_PIXELS = 50_000_000
MAX_IMAGE_DIMENSION = 20_000
MAX_IMAGE_PIXELS = 50_000_000
MAX_OCR_DIMENSION = 20_000
MAX_OCR_PIXELS = 50_000_000
MAX_ARCHIVE_MEMBERS = 10_000
MAX_ARCHIVE_MEMBER_BYTES = 100 * 1024 * 1024
MAX_ARCHIVE_UNCOMPRESSED_BYTES = 250 * 1024 * 1024
MAX_ARCHIVE_COMPRESSION_RATIO = 200
MIN_ARCHIVE_RATIO_CHECK_BYTES = 10 * 1024 * 1024
MAX_ARCHIVE_CENTRAL_DIRECTORY_BYTES = 32 * 1024 * 1024
SHEET_ROW_GROUP_ROWS = 25
JSON_INLINE_ITEMS = 25
JSON_INLINE_CHARACTERS = 1_000

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

_TEXT_MIMES = {"text/plain", "text/markdown", "text/x-markdown"}
_HTML_MIMES = {"text/html", "application/xhtml+xml"}
_CSV_MIMES = {"text/csv", "application/csv"}
_JSON_MIMES = {"application/json", "text/json"}
_IMAGE_MIMES = {
    "image/png",
    "image/jpeg",
    "image/tiff",
    "image/bmp",
    "image/webp",
    "image/gif",
}
_OCR_LANGUAGES = re.compile(r"^[A-Za-z0-9_+-]+$")


def extract_file(
    file_id: str,
    revision: str,
    path: Path,
    mime_type: str,
    *,
    ocr_languages: str = "eng",
) -> ExtractedDocument:
    """Extract a verified local artifact using its declared media type."""

    _validate_identity(file_id, revision)
    source = require_bounded_file(path, "payload")
    _validate_ocr_languages(ocr_languages)

    try:
        if mime_type == "application/pdf":
            if source.stat().st_size > MAX_PDF_FILE_BYTES:
                raise DriveRagError(
                    "PDF exceeds the 32 MiB extraction limit",
                    code="EXTRACTION_LIMIT_EXCEEDED",
                )
            return _extract_pdf(file_id, revision, source, ocr_languages)
        if mime_type == DOCX_MIME:
            return _extract_docx(file_id, revision, source)
        if mime_type == XLSX_MIME:
            return _extract_xlsx(file_id, revision, source)
        if mime_type == PPTX_MIME:
            return _extract_pptx(file_id, revision, source)
        if mime_type in _TEXT_MIMES:
            return _extract_text(file_id, revision, source, mime_type)
        if mime_type in _HTML_MIMES:
            return _extract_html(file_id, revision, source)
        if mime_type in _CSV_MIMES:
            return _extract_csv(file_id, revision, source)
        if mime_type in _JSON_MIMES:
            return _extract_json(file_id, revision, source)
        if mime_type in _IMAGE_MIMES:
            return _extract_image(file_id, revision, source, ocr_languages)
    except DriveRagError:
        raise
    except (
        OSError,
        ValueError,
        KeyError,
        TypeError,
        RuntimeError,
        json.JSONDecodeError,
        csv.Error,
        fitz.FileDataError,
        UnidentifiedImageError,
        BadZipFile,
        InvalidFileException,
    ) as exc:
        raise DriveRagError(
            f"could not extract {source.name}: {exc}", code="EXTRACTION_FAILED"
        ) from exc
    raise DriveRagError(
        f"unsupported indexing format: {mime_type}", code="UNSUPPORTED_FORMAT"
    )


def extract_native_structured(
    file_id: str, revision: str, raw: Mapping[str, object]
) -> ExtractedDocument:
    """Normalize exactly one connector-native structured schema."""

    _validate_identity(file_id, revision)
    if not isinstance(raw, Mapping):
        raise _structured_error("root must be an object")
    kind = raw.get("kind")
    if kind == "document":
        return _extract_native_document(file_id, revision, raw)
    if kind == "spreadsheet":
        return _extract_native_spreadsheet(file_id, revision, raw)
    if kind == "presentation":
        return _extract_native_presentation(file_id, revision, raw)
    raise _structured_error("kind must be document, spreadsheet, or presentation")


def require_bounded_file(path: Path, label: str) -> Path:
    """Require a regular, non-symlink file within the extraction byte limit."""

    source = Path(path)
    if source.is_symlink():
        raise DriveRagError(f"{label} must not be a symlink", code="UNSAFE_PATH")
    try:
        if not source.is_file():
            raise DriveRagError(
                f"{label} is not a readable file: {source}",
                code="EXTRACTION_FAILED",
            )
        size = source.stat().st_size
    except OSError as exc:
        raise DriveRagError(
            f"could not inspect {label} {source}: {exc}", code="EXTRACTION_FAILED"
        ) from exc
    if size > MAX_FILE_BYTES:
        raise DriveRagError(
            f"{label} exceeds the 250 MiB extraction limit",
            code="EXTRACTION_LIMIT_EXCEEDED",
        )
    return source


def _extract_native_document(
    file_id: str, revision: str, raw: Mapping[str, object]
) -> ExtractedDocument:
    _require_exact_keys(raw, {"kind", "sections"})
    sections = raw["sections"]
    if not isinstance(sections, list):
        raise _structured_error("document sections must be a list")
    if len(sections) > MAX_STRUCTURED_ITEMS:
        raise _structured_error("document exceeds the 20,000 section limit")
    blocks: list[ExtractedBlock] = []
    text_bytes = 0
    for section in sections:
        if not isinstance(section, Mapping):
            raise _structured_error("document section must be an object")
        _require_exact_keys(section, {"locator", "text"})
        locator = _nonempty_string(section["locator"], "section locator")
        text = _string(section["text"], "section text")
        if not locator.startswith("section:"):
            raise _structured_error("document locator must start with section:")
        text_bytes = _add_structured_bytes(text_bytes, locator, text)
        blocks.append(ExtractedBlock(locator, _normalize(text), {}))
    return ExtractedDocument(file_id, revision, tuple(blocks), "document")


def _extract_native_spreadsheet(
    file_id: str, revision: str, raw: Mapping[str, object]
) -> ExtractedDocument:
    _require_exact_keys(raw, {"kind", "sheets"})
    sheets = raw["sheets"]
    if not isinstance(sheets, list):
        raise _structured_error("spreadsheet sheets must be a list")
    if len(sheets) > MAX_STRUCTURED_ITEMS:
        raise _structured_error("spreadsheet exceeds the 20,000 sheet limit")
    blocks: list[ExtractedBlock] = []
    cells_seen = 0
    text_bytes = 0
    output_bytes = 0
    names: set[str] = set()
    for sheet in sheets:
        if not isinstance(sheet, Mapping):
            raise _structured_error("spreadsheet sheet must be an object")
        _require_exact_keys(sheet, {"name", "rows"})
        name = _nonempty_string(sheet["name"], "sheet name")
        if name in names:
            raise _structured_error("spreadsheet sheet names must be unique")
        names.add(name)
        text_bytes = _add_structured_bytes(text_bytes, name)
        rows = sheet["rows"]
        if not isinstance(rows, list) or any(not isinstance(row, list) for row in rows):
            raise _structured_error("spreadsheet rows must be lists")
        width = max((len(row) for row in rows), default=1)
        cells_seen += sum(max(1, len(row)) for row in rows)
        if cells_seen > MAX_SHEET_CELLS:
            raise _structured_error("spreadsheet exceeds the 50,000 cell limit")
        rendered: list[list[str]] = []
        for row in rows:
            rendered_row = [_render_cell(cell) for cell in row]
            text_bytes = _add_structured_bytes(text_bytes, *rendered_row)
            rendered.append(rendered_row)
        sheet_blocks = _sheet_blocks(name, rendered, include_range=True)
        output_bytes = _add_structured_bytes(
            output_bytes, *(block.text for block in sheet_blocks)
        )
        blocks.extend(sheet_blocks)
    return ExtractedDocument(file_id, revision, tuple(blocks), "spreadsheet")


def _extract_native_presentation(
    file_id: str, revision: str, raw: Mapping[str, object]
) -> ExtractedDocument:
    _require_exact_keys(raw, {"kind", "slides"})
    slides = raw["slides"]
    if not isinstance(slides, list):
        raise _structured_error("presentation slides must be a list")
    if len(slides) > MAX_STRUCTURED_ITEMS:
        raise _structured_error("presentation exceeds the 20,000 slide limit")
    blocks: list[ExtractedBlock] = []
    numbers: set[int] = set()
    text_bytes = 0
    for slide in slides:
        if not isinstance(slide, Mapping):
            raise _structured_error("presentation slide must be an object")
        _require_exact_keys(slide, {"number", "text", "notes"})
        number = slide["number"]
        if (
            not isinstance(number, int)
            or isinstance(number, bool)
            or number <= 0
            or number in numbers
        ):
            raise _structured_error("presentation slide number must be unique and positive")
        numbers.add(number)
        text = _string(slide["text"], "slide text")
        notes = _string(slide["notes"], "slide notes")
        text_bytes = _add_structured_bytes(text_bytes, text, notes)
        combined = _with_notes(text, notes)
        blocks.append(ExtractedBlock(f"slide:{number}", combined, {}))
    return ExtractedDocument(file_id, revision, tuple(blocks), "presentation")


def _extract_pdf(
    file_id: str, revision: str, path: Path, ocr_languages: str
) -> ExtractedDocument:
    blocks: list[ExtractedBlock] = []
    with fitz.open(path) as document:
        if document.needs_pass:
            raise DriveRagError(
                "encrypted PDF cannot be extracted", code="EXTRACTION_FAILED"
            )
        if document.page_count > MAX_PDF_PAGES:
            raise DriveRagError(
                f"PDF exceeds the {MAX_PDF_PAGES:,} page limit",
                code="EXTRACTION_LIMIT_EXCEEDED",
            )
        ocr_deadline = clock.monotonic() + OCR_TIMEOUT_SECONDS
        for index, page in enumerate(document):
            native_text = _normalize(page.get_text("text"))
            if len(native_text) < PDF_OCR_TEXT_THRESHOLD:
                scale = PDF_OCR_DPI / 72
                if (
                    not math.isfinite(page.rect.width)
                    or not math.isfinite(page.rect.height)
                    or page.rect.width <= 0
                    or page.rect.height <= 0
                ):
                    raise DriveRagError(
                        "PDF page exceeds the OCR pixel limit",
                        code="EXTRACTION_LIMIT_EXCEEDED",
                    )
                pixel_width = math.ceil(page.rect.width * scale)
                pixel_height = math.ceil(page.rect.height * scale)
                if pixel_width * pixel_height > MAX_PDF_OCR_PIXELS:
                    raise DriveRagError(
                        "PDF page exceeds the OCR pixel limit",
                        code="EXTRACTION_LIMIT_EXCEEDED",
                    )
                pixmap = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
                image = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
                ocr_text = _ocr_image(image, ocr_languages, deadline=ocr_deadline)
                text = ocr_text or native_text
            else:
                text = native_text
            blocks.append(ExtractedBlock(f"page:{index + 1}", text, {}))
    return ExtractedDocument(file_id, revision, tuple(blocks), "pdf")


def _extract_docx(file_id: str, revision: str, path: Path) -> ExtractedDocument:
    _preflight_office_archive(path)
    document = Document(path)
    blocks: list[ExtractedBlock] = []
    locator = "section:Document"
    lines: list[str] = []

    def flush() -> None:
        if any(line.strip() for line in lines):
            blocks.append(ExtractedBlock(locator, "\n".join(lines).strip(), {}))

    for item in document.iter_inner_content():
        if hasattr(item, "style") and getattr(item.style, "name", "").startswith("Heading"):
            flush()
            lines = []
            heading = _normalize(item.text)
            locator = f"section:{heading or 'Untitled'}"
            if heading:
                lines.append(heading)
        elif hasattr(item, "rows"):
            for row in item.rows:
                lines.append("\t".join(_normalize(cell.text) for cell in row.cells))
        else:
            text = _normalize(getattr(item, "text", ""))
            if text:
                lines.append(text)
    flush()
    return ExtractedDocument(file_id, revision, tuple(blocks), "docx")


def _extract_xlsx(file_id: str, revision: str, path: Path) -> ExtractedDocument:
    _preflight_office_archive(path)
    workbook = load_workbook(path, read_only=True, data_only=True, keep_links=False)
    blocks: list[ExtractedBlock] = []
    cells_seen = 0
    output_bytes = 0
    try:
        for sheet in workbook.worksheets:
            declared = sheet.max_row * sheet.max_column
            if declared > MAX_SHEET_CELLS or cells_seen + declared > MAX_SHEET_CELLS:
                raise DriveRagError(
                    "spreadsheet exceeds the 50,000 cell limit",
                    code="EXTRACTION_LIMIT_EXCEEDED",
                )
            rows = [
                [_render_office_cell(cell) for cell in row]
                for row in sheet.iter_rows(values_only=True)
            ]
            cells_seen += sum(len(row) for row in rows)
            if cells_seen > MAX_SHEET_CELLS:
                raise DriveRagError(
                    "spreadsheet exceeds the 50,000 cell limit",
                    code="EXTRACTION_LIMIT_EXCEEDED",
                )
            sheet_blocks = _sheet_blocks(sheet.title, rows, include_range=False)
            output_bytes = _add_structured_bytes(
                output_bytes, *(block.text for block in sheet_blocks)
            )
            blocks.extend(sheet_blocks)
    finally:
        workbook.close()
    return ExtractedDocument(file_id, revision, tuple(blocks), "xlsx")


def _extract_pptx(file_id: str, revision: str, path: Path) -> ExtractedDocument:
    _preflight_office_archive(path)
    presentation = Presentation(path)
    blocks: list[ExtractedBlock] = []
    for number, slide in enumerate(presentation.slides, start=1):
        texts = [
            _normalize(shape.text)
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and _normalize(shape.text)
        ]
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = _normalize(slide.notes_slide.notes_text_frame.text)
        blocks.append(
            ExtractedBlock(f"slide:{number}", _with_notes("\n".join(texts), notes), {})
        )
    return ExtractedDocument(file_id, revision, tuple(blocks), "pptx")


def _extract_text(
    file_id: str, revision: str, path: Path, mime_type: str
) -> ExtractedDocument:
    text = _read_text(path)
    if mime_type in {"text/markdown", "text/x-markdown"}:
        blocks = _heading_blocks(text)
        document_format = "markdown"
    else:
        blocks = (ExtractedBlock("text:1", _normalize(text), {}),)
        document_format = "text"
    return ExtractedDocument(file_id, revision, blocks, document_format)


def _extract_html(file_id: str, revision: str, path: Path) -> ExtractedDocument:
    soup = BeautifulSoup(_read_text(path), "html.parser")
    for tag in soup(["script", "style", "noscript", "template"]):
        tag.decompose()
    lines: list[str] = []
    for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "tr"]):
        text = _normalize(tag.get_text(" ", strip=True))
        if not text:
            continue
        if tag.name and tag.name.startswith("h"):
            lines.append(f"# {text}")
        else:
            lines.append(text)
    blocks = _heading_blocks("\n\n".join(lines))
    if not blocks:
        visible = _normalize(soup.get_text(" ", strip=True))
        if visible:
            blocks = (ExtractedBlock("section:Document", visible, {}),)
    return ExtractedDocument(file_id, revision, blocks, "html")


def _extract_csv(file_id: str, revision: str, path: Path) -> ExtractedDocument:
    with path.open("r", encoding="utf-8-sig", newline="") as stream:
        rows = list(csv.reader(stream))
    cell_count = sum(len(row) for row in rows)
    if cell_count > MAX_SHEET_CELLS:
        raise DriveRagError(
            "CSV exceeds the 50,000 cell limit", code="EXTRACTION_LIMIT_EXCEEDED"
        )
    blocks = _sheet_blocks("CSV", rows, include_range=True)
    return ExtractedDocument(file_id, revision, blocks, "csv")


def _extract_json(file_id: str, revision: str, path: Path) -> ExtractedDocument:
    try:
        raw = json.loads(_read_text(path))
    except (json.JSONDecodeError, RecursionError) as exc:
        raise DriveRagError("invalid JSON payload", code="EXTRACTION_FAILED") from exc
    if _count_json_nodes(raw) > MAX_JSON_NODES:
        raise DriveRagError(
            "JSON exceeds the 20,000 node limit", code="EXTRACTION_LIMIT_EXCEEDED"
        )
    return ExtractedDocument(file_id, revision, _json_blocks(raw), "json")


def _extract_image(
    file_id: str, revision: str, path: Path, ocr_languages: str
) -> ExtractedDocument:
    try:
        with warnings.catch_warnings():
            warnings.simplefilter("error", Image.DecompressionBombWarning)
            with Image.open(path) as image:
                _validate_image_extent(
                    image.width,
                    image.height,
                    label="image dimension",
                    max_dimension=MAX_IMAGE_DIMENSION,
                    max_pixels=MAX_IMAGE_PIXELS,
                )
                image.load()
                text = _ocr_image(image, ocr_languages)
    except (Image.DecompressionBombError, Image.DecompressionBombWarning) as exc:
        raise DriveRagError(
            "image dimension exceeds the safe pixel limit",
            code="EXTRACTION_LIMIT_EXCEEDED",
        ) from exc
    return ExtractedDocument(
        file_id, revision, (ExtractedBlock("image:1", text, {}),), "image"
    )


def _ocr_image(
    image: Image.Image, languages: str, *, deadline: float | None = None
) -> str:
    _validate_image_extent(
        image.width,
        image.height,
        label="transformed OCR image",
        max_dimension=MAX_OCR_DIMENSION,
        max_pixels=MAX_OCR_PIXELS,
    )
    prepared = ImageOps.autocontrast(ImageOps.grayscale(image))
    ink_bounds = ImageOps.invert(prepared).getbbox()
    if ink_bounds is not None:
        left, top, right, bottom = ink_bounds
        padding = max(8, (bottom - top) * 2)
        prepared = prepared.crop(
            (
                max(0, left - padding),
                max(0, top - padding),
                min(prepared.width, right + padding),
                min(prepared.height, bottom + padding),
            )
        )
        _validate_image_extent(
            prepared.width,
            prepared.height,
            label="transformed OCR crop",
            max_dimension=MAX_OCR_DIMENSION,
            max_pixels=MAX_OCR_PIXELS,
        )
        ink_height = max(1, bottom - top)
        scale = max(1, (40 + ink_height - 1) // ink_height)
        if scale > 1:
            scaled_width = prepared.width * scale
            scaled_height = prepared.height * scale
            _validate_image_extent(
                scaled_width,
                scaled_height,
                label="transformed OCR resize",
                max_dimension=MAX_OCR_DIMENSION,
                max_pixels=MAX_OCR_PIXELS,
            )
            prepared = prepared.resize(
                (scaled_width, scaled_height),
                Image.Resampling.BICUBIC,
            )
        border = max(20, 10 * scale)
        _validate_image_extent(
            prepared.width + 2 * border,
            prepared.height + 2 * border,
            label="transformed OCR border",
            max_dimension=MAX_OCR_DIMENSION,
            max_pixels=MAX_OCR_PIXELS,
        )
        prepared = ImageOps.expand(prepared, border=border, fill="white")
    try:
        if deadline is None:
            deadline = clock.monotonic() + OCR_TIMEOUT_SECONDS
        timeout = _remaining_ocr_time(deadline)
        text = _normalize(
            pytesseract.image_to_string(
                prepared, lang=languages, config="--psm 6", timeout=timeout
            )
        )
        if not text:
            timeout = _remaining_ocr_time(deadline)
            text = _normalize(
                pytesseract.image_to_string(
                    prepared,
                    lang=languages,
                    config="--psm 7",
                    timeout=timeout,
                )
            )
        return text
    except RuntimeError as exc:
        raise DriveRagError("OCR timed out or failed", code="OCR_FAILED") from exc


def _validate_image_extent(
    width: object,
    height: object,
    *,
    label: str,
    max_dimension: int,
    max_pixels: int,
) -> None:
    if (
        not isinstance(width, int)
        or isinstance(width, bool)
        or not isinstance(height, int)
        or isinstance(height, bool)
        or width <= 0
        or height <= 0
        or width > max_dimension
        or height > max_dimension
        or width * height > max_pixels
    ):
        raise DriveRagError(
            f"{label} exceeds the safe dimension or pixel limit",
            code="EXTRACTION_LIMIT_EXCEEDED",
        )


def _preflight_office_archive(path: Path) -> None:
    _preflight_archive_directory(path)
    try:
        with ZipFile(path, "r") as archive:
            members = archive.infolist()
            if len(members) > MAX_ARCHIVE_MEMBERS:
                raise DriveRagError(
                    "Office archive exceeds the member count limit",
                    code="EXTRACTION_LIMIT_EXCEEDED",
                )
            names: set[str] = set()
            aggregate_size = 0
            for member in members:
                _validate_archive_member_path(member.filename)
                if member.filename in names:
                    raise DriveRagError(
                        "Office archive contains duplicate member paths",
                        code="EXTRACTION_FAILED",
                    )
                names.add(member.filename)
                _validate_archive_member_type(member)
                if member.flag_bits & 0x1:
                    raise DriveRagError(
                        "Office archive contains an encrypted member",
                        code="EXTRACTION_FAILED",
                    )
                if member.file_size < 0 or member.compress_size < 0:
                    raise DriveRagError(
                        "Office archive has invalid member sizes",
                        code="EXTRACTION_FAILED",
                    )
                if member.file_size > MAX_ARCHIVE_MEMBER_BYTES:
                    raise DriveRagError(
                        "Office archive member expansion exceeds the limit",
                        code="EXTRACTION_LIMIT_EXCEEDED",
                    )
                aggregate_size += member.file_size
                if aggregate_size > MAX_ARCHIVE_UNCOMPRESSED_BYTES:
                    raise DriveRagError(
                        "Office archive aggregate expansion exceeds the limit",
                        code="EXTRACTION_LIMIT_EXCEEDED",
                    )
                if member.file_size >= MIN_ARCHIVE_RATIO_CHECK_BYTES and (
                    member.compress_size == 0
                    or member.file_size / member.compress_size
                    > MAX_ARCHIVE_COMPRESSION_RATIO
                ):
                    raise DriveRagError(
                        "Office archive member compression ratio exceeds the limit",
                        code="EXTRACTION_LIMIT_EXCEEDED",
                    )
    except DriveRagError:
        raise
    except (BadZipFile, OSError, ValueError) as exc:
        raise DriveRagError(
            "Office payload is not a readable ZIP archive",
            code="EXTRACTION_FAILED",
        ) from exc


def _preflight_archive_directory(path: Path) -> None:
    """Bound central-directory work before ``ZipFile`` materializes entries."""

    try:
        file_size = path.stat().st_size
        if file_size < 22:
            raise ValueError("missing ZIP end record")
        tail_size = min(file_size, 22 + 65_535)
        with path.open("rb") as archive:
            archive.seek(file_size - tail_size)
            tail = archive.read(tail_size)
            eocd_index = _find_eocd(tail)
            if eocd_index is None:
                raise ValueError("missing ZIP end record")
            eocd_offset = file_size - tail_size + eocd_index
            (
                _signature,
                disk_number,
                directory_disk,
                disk_entries,
                total_entries,
                directory_size,
                directory_offset,
                _comment_size,
            ) = struct.unpack_from("<4s4H2LH", tail, eocd_index)

            uses_zip64 = (
                disk_entries == 0xFFFF
                or total_entries == 0xFFFF
                or directory_size == 0xFFFFFFFF
                or directory_offset == 0xFFFFFFFF
            )
            directory_end_limit = eocd_offset
            if uses_zip64:
                locator_offset = eocd_offset - 20
                if locator_offset < 0:
                    raise ValueError("missing ZIP64 locator")
                archive.seek(locator_offset)
                locator = archive.read(20)
                if len(locator) != 20:
                    raise ValueError("truncated ZIP64 locator")
                (
                    locator_signature,
                    zip64_disk,
                    zip64_offset,
                    total_disks,
                ) = struct.unpack("<4sLQL", locator)
                if locator_signature != b"PK\x06\x07":
                    raise ValueError("missing ZIP64 locator")
                if zip64_disk != 0 or total_disks != 1:
                    raise ValueError("multi-disk ZIP archives are unsupported")
                archive.seek(zip64_offset)
                zip64_record = archive.read(56)
                if len(zip64_record) != 56:
                    raise ValueError("truncated ZIP64 end record")
                (
                    zip64_signature,
                    zip64_record_size,
                    _version_made,
                    _version_needed,
                    disk_number,
                    directory_disk,
                    disk_entries,
                    total_entries,
                    directory_size,
                    directory_offset,
                ) = struct.unpack("<4sQ2H2L4Q", zip64_record)
                if zip64_signature != b"PK\x06\x06" or zip64_record_size < 44:
                    raise ValueError("invalid ZIP64 end record")
                if zip64_offset + 12 + zip64_record_size > locator_offset:
                    raise ValueError("invalid ZIP64 end record extent")
                directory_end_limit = zip64_offset

        if disk_number != 0 or directory_disk != 0 or disk_entries != total_entries:
            raise ValueError("multi-disk ZIP archives are unsupported")
        if total_entries > MAX_ARCHIVE_MEMBERS:
            raise DriveRagError(
                "Office archive exceeds the member count limit",
                code="EXTRACTION_LIMIT_EXCEEDED",
            )
        if directory_size > MAX_ARCHIVE_CENTRAL_DIRECTORY_BYTES:
            raise DriveRagError(
                "Office archive central directory exceeds the limit",
                code="EXTRACTION_LIMIT_EXCEEDED",
            )
        if (
            directory_offset > directory_end_limit
            or directory_size > directory_end_limit - directory_offset
        ):
            raise ValueError("invalid ZIP central directory extent")
        actual_entries = _count_central_directory_entries(
            path, directory_offset, directory_size
        )
        if actual_entries != total_entries:
            raise ValueError("inconsistent ZIP central directory entry count")
    except DriveRagError:
        raise
    except (OSError, ValueError, struct.error) as exc:
        raise DriveRagError(
            "Office payload is not a readable ZIP archive",
            code="EXTRACTION_FAILED",
        ) from exc


def _count_central_directory_entries(
    path: Path, directory_offset: int, directory_size: int
) -> int:
    with path.open("rb") as archive:
        archive.seek(directory_offset)
        directory = archive.read(directory_size)
    if len(directory) != directory_size:
        raise ValueError("truncated ZIP central directory")

    entry_count = 0
    position = 0
    while position < directory_size:
        if position + 46 > directory_size:
            raise ValueError("truncated ZIP central directory entry")
        if directory[position : position + 4] != b"PK\x01\x02":
            raise ValueError("invalid ZIP central directory entry")
        name_size, extra_size, comment_size = struct.unpack_from(
            "<3H", directory, position + 28
        )
        entry_size = 46 + name_size + extra_size + comment_size
        if entry_size > directory_size - position:
            raise ValueError("truncated ZIP central directory entry")
        entry_count += 1
        if entry_count > MAX_ARCHIVE_MEMBERS:
            raise DriveRagError(
                "Office archive exceeds the member count limit",
                code="EXTRACTION_LIMIT_EXCEEDED",
            )
        position += entry_size
    return entry_count


def _find_eocd(tail: bytes) -> int | None:
    search_end = len(tail)
    while True:
        index = tail.rfind(b"PK\x05\x06", 0, search_end)
        if index < 0:
            return None
        if index + 22 <= len(tail):
            comment_size = struct.unpack_from("<H", tail, index + 20)[0]
            if index + 22 + comment_size == len(tail):
                return index
        search_end = index


def _validate_archive_member_path(name: object) -> None:
    if not isinstance(name, str) or not name or "\\" in name:
        raise DriveRagError(
            "unsafe archive member path", code="EXTRACTION_FAILED"
        )
    trimmed = name[:-1] if name.endswith("/") else name
    raw_parts = trimmed.split("/")
    path = PurePosixPath(name)
    if (
        not trimmed
        or path.is_absolute()
        or name.startswith("/")
        or any(part in {"", ".", ".."} for part in raw_parts)
        or (raw_parts and raw_parts[0].endswith(":"))
        or any(unicodedata.category(character).startswith("C") for character in name)
    ):
        raise DriveRagError(
            "unsafe archive member path", code="EXTRACTION_FAILED"
        )


def _validate_archive_member_type(member: object) -> None:
    external_attr = getattr(member, "external_attr", 0)
    create_system = getattr(member, "create_system", 0)
    is_directory = member.is_dir()  # type: ignore[union-attr]
    mode = (external_attr >> 16) & 0xFFFF
    file_type = stat.S_IFMT(mode)
    if create_system == 3 and file_type not in {
        0,
        stat.S_IFDIR if is_directory else stat.S_IFREG,
    }:
        raise DriveRagError(
            "unsafe archive member type", code="EXTRACTION_FAILED"
        )


def _heading_blocks(text: str) -> tuple[ExtractedBlock, ...]:
    blocks: list[ExtractedBlock] = []
    locator = "section:Document"
    lines: list[str] = []
    for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
        heading = re.match(r"^\s{0,3}#{1,6}\s+(.+?)\s*#*\s*$", line)
        if heading:
            if any(item.strip() for item in lines):
                blocks.append(ExtractedBlock(locator, _normalize("\n".join(lines)), {}))
            title = _normalize(heading.group(1))
            locator = f"section:{title or 'Untitled'}"
            lines = [title] if title else []
        elif line.strip():
            lines.append(line.strip())
    if any(item.strip() for item in lines):
        blocks.append(ExtractedBlock(locator, _normalize("\n".join(lines)), {}))
    return tuple(blocks)


def _count_json_nodes(raw: object) -> int:
    count = 0
    stack = [raw]
    while stack:
        node = stack.pop()
        count += 1
        if count > MAX_JSON_NODES:
            return count
        if isinstance(node, dict):
            stack.extend(node.values())
        elif isinstance(node, list):
            stack.extend(node)
    return count


def _read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8-sig")
    except UnicodeDecodeError as exc:
        raise DriveRagError(
            f"payload is not valid UTF-8: {path.name}", code="EXTRACTION_FAILED"
        ) from exc


def _column_name(index: int) -> str:
    if index <= 0:
        return "A"
    result = ""
    while index:
        index, remainder = divmod(index - 1, 26)
        result = chr(65 + remainder) + result
    return result


def _sheet_blocks(
    name: str, rows: list[list[str]], *, include_range: bool
) -> tuple[ExtractedBlock, ...]:
    width = max((len(row) for row in rows), default=1)
    header = rows[0] if rows else []
    data_rows = rows[1:] if rows else []
    groups = [
        data_rows[index : index + SHEET_ROW_GROUP_ROWS]
        for index in range(0, len(data_rows), SHEET_ROW_GROUP_ROWS)
    ] or [[]]
    blocks: list[ExtractedBlock] = []
    output_bytes = 0
    for index, group in enumerate(groups):
        source_start = 2 + index * SHEET_ROW_GROUP_ROWS
        source_end = source_start + len(group) - 1
        if index == 0:
            locator_start = 1
            locator_end = max(1, source_end)
        else:
            locator_start = source_start
            locator_end = source_end
        range_locator = (
            f"sheet:{name}!A{locator_start}:"
            f"{_column_name(width)}{locator_end}"
        )
        content_rows = ([header] if rows else []) + group
        text = "\n".join("\t".join(row).rstrip() for row in content_rows).strip()
        output_bytes = _add_structured_bytes(output_bytes, text)
        blocks.append(
            ExtractedBlock(
                range_locator,
                text,
                {
                    "sheet_name": name,
                    "header_row": 1 if rows else None,
                    "row_start": locator_start,
                    "row_end": locator_end,
                },
            )
        )
    if not include_range and len(blocks) == 1:
        block = blocks[0]
        blocks[0] = ExtractedBlock(f"sheet:{name}", block.text, block.metadata)
    return tuple(blocks)


def _json_blocks(raw: object) -> tuple[ExtractedBlock, ...]:
    blocks: list[ExtractedBlock] = []
    pending: list[tuple[str, object]] = [("$", raw)]
    while pending:
        path, value = pending.pop()
        if isinstance(value, dict):
            items = sorted(value.items(), key=lambda item: item[0])
            rendered = json.dumps(value, ensure_ascii=False, sort_keys=True)
            if not items or (
                len(items) <= JSON_INLINE_ITEMS
                and all(not isinstance(child, (dict, list)) for _, child in items)
                and len(rendered) <= JSON_INLINE_CHARACTERS
            ):
                blocks.append(ExtractedBlock(f"json:{path}", rendered, {}))
            else:
                for key, child in reversed(items):
                    pending.append((f"{path}{_json_key(key)}", child))
        elif isinstance(value, list):
            rendered = json.dumps(value, ensure_ascii=False, sort_keys=True)
            if not value or (
                len(value) <= JSON_INLINE_ITEMS
                and all(not isinstance(child, (dict, list)) for child in value)
                and len(rendered) <= JSON_INLINE_CHARACTERS
            ):
                blocks.append(ExtractedBlock(f"json:{path}", rendered, {}))
            else:
                for index in range(len(value) - 1, -1, -1):
                    pending.append((f"{path}[{index}]", value[index]))
        else:
            text = value if isinstance(value, str) else json.dumps(value, ensure_ascii=False)
            blocks.append(ExtractedBlock(f"json:{path}", str(text), {}))
    return tuple(blocks)


def _json_key(key: str) -> str:
    if re.fullmatch(r"[A-Za-z_][A-Za-z0-9_]*", key):
        return f".{key}"
    return f"[{json.dumps(key, ensure_ascii=False)}]"


def _render_cell(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, (str, int, float)):
        return str(value)
    raise _structured_error("spreadsheet cells must be scalar values")


def _render_office_cell(value: object) -> str:
    if isinstance(value, (datetime, date, time)):
        return value.isoformat()
    return _render_cell(value)


def _add_structured_bytes(total: int, *values: str) -> int:
    total += sum(len(value.encode("utf-8")) for value in values)
    if total > MAX_FILE_BYTES:
        raise _structured_error("content exceeds the 250 MiB extraction limit")
    return total


def _remaining_ocr_time(deadline: float) -> float:
    remaining = deadline - clock.monotonic()
    if remaining <= 0:
        raise DriveRagError("OCR exceeded the 10 minute budget", code="OCR_FAILED")
    return min(float(OCR_TIMEOUT_SECONDS), remaining)


def _with_notes(text: str, notes: str) -> str:
    visible = _normalize(text)
    speaker_notes = _normalize(notes)
    if visible and speaker_notes:
        return f"{visible}\nNotes:\n{speaker_notes}"
    return visible or speaker_notes


def _normalize(text: str) -> str:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    normalized = re.sub(r"[ \t]+", " ", normalized)
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    return normalized.strip()


def _validate_identity(file_id: object, revision: object) -> None:
    if not isinstance(file_id, str) or not file_id.strip():
        raise DriveRagError("file_id must not be empty", code="INVALID_EXTRACTION")
    if not isinstance(revision, str) or not revision.strip():
        raise DriveRagError("revision must not be empty", code="INVALID_EXTRACTION")


def _validate_ocr_languages(languages: str) -> None:
    if not isinstance(languages, str) or not _OCR_LANGUAGES.fullmatch(languages):
        raise DriveRagError("invalid OCR languages", code="INVALID_ARGUMENTS")


def _require_exact_keys(raw: Mapping[str, object], expected: set[str]) -> None:
    if set(raw) != expected:
        raise _structured_error("object fields do not match the structured schema")


def _string(value: object, field: str) -> str:
    if not isinstance(value, str):
        raise _structured_error(f"{field} must be a string")
    return value


def _nonempty_string(value: object, field: str) -> str:
    result = _string(value, field)
    if not result.strip():
        raise _structured_error(f"{field} must not be empty")
    return result


def _structured_error(message: str) -> DriveRagError:
    return DriveRagError(
        f"invalid structured content: {message}", code="INVALID_STRUCTURED_CONTENT"
    )
