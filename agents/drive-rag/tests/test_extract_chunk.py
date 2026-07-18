import json
from datetime import date
from pathlib import Path
import stat
import struct
import subprocess
import sys
import zipfile
import zlib

import fitz
from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont, PngImagePlugin
from pptx import Presentation
import pytest
import drive_rag_lib.extract as extract_module

from drive_rag_lib.chunk import chunk_document
from drive_rag_lib.extract import (
    extract_file,
    extract_native_structured,
    require_bounded_file,
)
from drive_rag_lib.models import ExtractedBlock, ExtractedDocument
from drive_rag_lib.protocol import DriveRagError, read_json


DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class WordCounter:
    def count(self, text):
        return len(text.split())


class CharacterCounter:
    def count(self, text):
        return len(text)


def test_office_locators(tmp_path: Path):
    docx_path = tmp_path / "source.docx"
    document = Document()
    document.add_heading("Policy", 1)
    document.add_paragraph("Retention is thirty days.")
    document.save(docx_path)

    xlsx_path = tmp_path / "source.xlsx"
    workbook = Workbook()
    workbook.active.title = "Rates"
    workbook.active.append(["Region", "Rate"])
    workbook.active.append(["EU", 20])
    workbook.save(xlsx_path)

    pptx_path = tmp_path / "source.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Launch"
    slide.placeholders[1].text = "Ship safely"
    presentation.save(pptx_path)

    assert "Policy" in extract_file("doc", "1", docx_path, DOCX_MIME).blocks[0].text
    assert extract_file("sheet", "1", xlsx_path, XLSX_MIME).blocks[0].locator == "sheet:Rates"
    assert extract_file("slides", "1", pptx_path, PPTX_MIME).blocks[0].locator == "slide:1"


def test_image_ocr(tmp_path: Path):
    path = tmp_path / "scan.png"
    image = Image.new("RGB", (1200, 300), "white")
    ImageDraw.Draw(image).text(
        (40, 100), "DRIVE OCR 2026", fill="black", stroke_width=1
    )
    image.save(path)

    result = extract_file("image", "1", path, "image/png", ocr_languages="eng")

    assert "2026" in " ".join(block.text for block in result.blocks)


def test_pdf_byte_limit_fails_before_opening_document(tmp_path: Path, monkeypatch):
    path = tmp_path / "oversized.pdf"
    path.write_bytes(b"not-a-pdf")
    monkeypatch.setattr(extract_module, "MAX_PDF_FILE_BYTES", 1)

    with pytest.raises(DriveRagError) as raised:
        extract_file("pdf", "1", path, "application/pdf")

    assert raised.value.code == "EXTRACTION_LIMIT_EXCEEDED"


def test_scanned_pdf_uses_ocr(tmp_path: Path):
    image = Image.new("RGB", (1600, 400), "white")
    font = ImageFont.truetype(
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 52
    )
    ImageDraw.Draw(image).text(
        (50, 130), "SCANNED POLICY 2026", fill="black", font=font
    )
    scan = tmp_path / "scan.png"
    image.save(scan)
    pdf_path = tmp_path / "scan.pdf"
    pdf = fitz.open()
    page = pdf.new_page(width=800, height=200)
    page.insert_image(page.rect, filename=str(scan))
    pdf.save(pdf_path)
    pdf.close()

    result = extract_file("pdf", "1", pdf_path, "application/pdf")

    assert result.blocks[0].locator == "page:1"
    assert "2026" in result.blocks[0].text


def test_image_dimensions_are_rejected_before_pixel_decode(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    path = tmp_path / "oversized.png"
    Image.new("RGB", (20, 20), "white").save(path)
    sentinel = tmp_path / "previous.txt"
    sentinel.write_text("previous-version-active", encoding="utf-8")
    monkeypatch.setattr(extract_module, "MAX_IMAGE_DIMENSION", 10, raising=False)

    def forbidden_load(*args, **kwargs):
        raise AssertionError("pixel data was decoded")

    monkeypatch.setattr(PngImagePlugin.PngImageFile, "load", forbidden_load)

    with pytest.raises(DriveRagError, match="dimension") as error:
        extract_file("image", "1", path, "image/png")
    assert error.value.code == "EXTRACTION_LIMIT_EXCEEDED"
    assert sentinel.read_text(encoding="utf-8") == "previous-version-active"


def test_ocr_rejects_unsafe_transformed_dimensions_before_resize(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    path = tmp_path / "wide.png"
    image = Image.new("L", (10_000, 20), "white")
    image.putpixel((0, 10), 0)
    image.putpixel((9_999, 10), 0)
    image.save(path)

    def forbidden_resize(*args, **kwargs):
        raise AssertionError("unsafe transformed image was resized")

    monkeypatch.setattr(Image.Image, "resize", forbidden_resize)

    with pytest.raises(DriveRagError, match="transformed") as error:
        extract_file("image", "1", path, "image/png")
    assert error.value.code == "EXTRACTION_LIMIT_EXCEEDED"


def test_decompression_bomb_header_returns_typed_error_without_content_leak(
    tmp_path: Path,
):
    def png_chunk(kind: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + kind
            + data
            + struct.pack(">I", zlib.crc32(kind + data) & 0xFFFFFFFF)
        )

    path = tmp_path / "bomb.png"
    header = struct.pack(">IIBBBBB", 100_000, 100_000, 8, 2, 0, 0, 0)
    path.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        + png_chunk(b"IHDR", header)
        + png_chunk(b"IEND", b"")
    )
    sentinel = tmp_path / "previous.txt"
    sentinel.write_text("SECRET_DOCUMENT_BODY", encoding="utf-8")

    with pytest.raises(DriveRagError) as error:
        extract_file("image", "1", path, "image/png")

    assert error.value.code == "EXTRACTION_LIMIT_EXCEEDED"
    assert "SECRET_DOCUMENT_BODY" not in str(error.value)
    assert sentinel.read_text(encoding="utf-8") == "SECRET_DOCUMENT_BODY"


def test_native_schemas_and_stable_chunk_ids():
    sheet = {
        "kind": "spreadsheet",
        "sheets": [
            {
                "name": "Budget",
                "rows": [["Item", "Cost"], ["Hosting", "30"]],
            }
        ],
    }
    document = {
        "kind": "document",
        "sections": [{"locator": "section:Policy", "text": "Retention policy"}],
    }
    presentation = {
        "kind": "presentation",
        "slides": [{"number": 1, "text": "Launch", "notes": "Internal"}],
    }

    assert extract_native_structured("sheet", "7", sheet).blocks[0].locator == "sheet:Budget!A1:B2"
    assert extract_native_structured("doc", "7", document).blocks[0].locator == "section:Policy"
    assert extract_native_structured("slides", "7", presentation).blocks[0].locator == "slide:1"

    source = ExtractedDocument(
        "file-a",
        "1",
        (ExtractedBlock("section:One", "one two three four five " * 10, {}),),
        "text",
    )
    first = chunk_document(source, WordCounter(), 20, 4)
    second = chunk_document(source, WordCounter(), 20, 4)

    assert [item.chunk_id for item in first] == [item.chunk_id for item in second]
    assert max(WordCounter().count(item.text) for item in first) <= 20
    assert len(first) > 1


@pytest.mark.parametrize(
    "raw",
    [
        {"kind": "document", "sections": [], "extra": True},
        {"kind": "document", "sections": [{"locator": "", "text": "x"}]},
        {"kind": "spreadsheet", "sheets": [{"name": "S", "rows": [[object()]]}]},
        {"kind": "presentation", "slides": [{"number": True, "text": "x", "notes": ""}]},
    ],
)
def test_native_schema_is_strict(raw):
    with pytest.raises(DriveRagError, match="structured"):
        extract_native_structured("file-a", "1", raw)


def test_native_sheet_limit_is_enforced():
    raw = {
        "kind": "spreadsheet",
        "sheets": [{"name": "Large", "rows": [["x"]] * 50_001}],
    }
    with pytest.raises(DriveRagError, match="50,000"):
        extract_native_structured("file-a", "1", raw)


def test_native_sheet_bounds_empty_sheets_and_rows():
    too_many_sheets = {
        "kind": "spreadsheet",
        "sheets": [
            {"name": f"sheet-{index}", "rows": []} for index in range(20_001)
        ],
    }
    too_many_empty_rows = {
        "kind": "spreadsheet",
        "sheets": [{"name": "Empty", "rows": [[] for _ in range(50_001)]}],
    }

    with pytest.raises(DriveRagError, match="20,000"):
        extract_native_structured("file-a", "1", too_many_sheets)
    with pytest.raises(DriveRagError, match="50,000"):
        extract_native_structured("file-a", "1", too_many_empty_rows)


def test_sheet_header_repetition_cannot_amplify_past_output_limit(
    monkeypatch: pytest.MonkeyPatch,
):
    monkeypatch.setattr(extract_module, "MAX_FILE_BYTES", 100)
    raw = {
        "kind": "spreadsheet",
        "sheets": [
            {
                "name": "Budget",
                "rows": [["H" * 40]] + [[str(index)] for index in range(30)],
            }
        ],
    }

    with pytest.raises(DriveRagError, match="250 MiB"):
        extract_native_structured("file-a", "1", raw)


def test_native_docs_and_slides_have_item_limits():
    document = {
        "kind": "document",
        "sections": [
            {"locator": f"section:{index}", "text": "x"}
            for index in range(20_001)
        ],
    }
    presentation = {
        "kind": "presentation",
        "slides": [
            {"number": index + 1, "text": "x", "notes": ""}
            for index in range(20_001)
        ],
    }

    for raw in (document, presentation):
        with pytest.raises(DriveRagError, match="20,000"):
            extract_native_structured("file-a", "1", raw)


def test_bounded_file_guard_rejects_sparse_structured_input(tmp_path: Path):
    path = tmp_path / "large.structured.json"
    with path.open("wb") as stream:
        stream.truncate(250 * 1024 * 1024 + 1)

    with pytest.raises(DriveRagError, match="250 MiB"):
        require_bounded_file(path, "structured content")


def test_native_sheet_uses_bounded_row_groups_with_header_context():
    rows = [["Item", "Cost"]] + [[f"item-{index}", index] for index in range(60)]
    raw = {"kind": "spreadsheet", "sheets": [{"name": "Budget", "rows": rows}]}

    result = extract_native_structured("sheet", "7", raw)

    assert len(result.blocks) == 3
    assert all(block.text.splitlines()[0] == "Item\tCost" for block in result.blocks)
    assert [block.locator for block in result.blocks] == [
        "sheet:Budget!A1:B26",
        "sheet:Budget!A27:B51",
        "sheet:Budget!A52:B61",
    ]


def test_sheet_chunks_repeat_header_and_keep_rows_together():
    rows = [["Item", "Cost"]] + [
        [f"hosting-{index}", f"cost-{index}"] for index in range(12)
    ]
    document = extract_native_structured(
        "sheet", "7", {"kind": "spreadsheet", "sheets": [{"name": "Budget", "rows": rows}]}
    )

    chunks = chunk_document(document, WordCounter(), 8, 2)

    assert len(chunks) > 1
    assert all(chunk.text.splitlines()[0] == "Item\tCost" for chunk in chunks)
    source_rows = {"\t".join(row) for row in rows[1:]}
    chunk_rows = {
        line
        for chunk in chunks
        for line in chunk.text.splitlines()[1:]
        if line
    }
    assert chunk_rows == source_rows


def test_oversized_sheet_row_accounts_for_header_separator_tokens():
    document = ExtractedDocument(
        "sheet",
        "1",
        (
            ExtractedBlock(
                "sheet:Budget!A1:A2", "H\nabcdefghijkl", {"sheet_name": "Budget"}
            ),
        ),
        "spreadsheet",
    )

    chunks = chunk_document(document, CharacterCounter(), 5, 1)

    assert chunks
    assert all(chunk.text.startswith("H\n") for chunk in chunks)
    assert all(CharacterCounter().count(chunk.text) <= 5 for chunk in chunks)
    reconstructed = chunks[0].text[2:]
    for chunk in chunks[1:]:
        reconstructed += chunk.text[3:]
    assert reconstructed == "abcdefghijkl"


def test_nested_json_preserves_object_array_boundaries(tmp_path: Path):
    path = tmp_path / "nested.json"
    path.write_text(
        json.dumps(
            {
                "departments": [
                    {"name": "Alpha", "budget": 10},
                    {"name": "Beta", "budget": 20},
                ]
            }
        ),
        encoding="utf-8",
    )

    result = extract_file("json", "1", path, "application/json")

    assert [block.locator for block in result.blocks] == [
        "json:$.departments[0]",
        "json:$.departments[1]",
    ]
    assert "Alpha" in result.blocks[0].text
    assert "Beta" in result.blocks[1].text


def test_pdf_rejects_oversized_200_dpi_raster_before_allocation(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    path = tmp_path / "huge-page.pdf"
    pdf = fitz.open()
    pdf.new_page(width=20_000, height=20_000)
    pdf.save(path)
    pdf.close()

    def forbidden_render(*args, **kwargs):
        raise AssertionError("oversized page was rendered")

    monkeypatch.setattr(fitz.Page, "get_pixmap", forbidden_render)

    with pytest.raises(DriveRagError, match="pixel"):
        extract_file("pdf", "1", path, "application/pdf")


def test_text_html_csv_json_and_pdf_are_structured(tmp_path: Path):
    markdown = tmp_path / "guide.md"
    markdown.write_text("# Setup\n\nInstall it.\n\n## Verify\n\nRun tests.", encoding="utf-8")
    html = tmp_path / "page.html"
    html.write_text(
        "<html><style>SECRET_STYLE</style><script>SECRET_SCRIPT</script>"
        "<h1>Visible</h1><p>Body</p></html>",
        encoding="utf-8",
    )
    csv_path = tmp_path / "rates.csv"
    csv_path.write_text("Region,Rate\nEU,20\n", encoding="utf-8")
    json_path = tmp_path / "data.json"
    json_path.write_text(json.dumps({"policy": {"days": 30}}), encoding="utf-8")
    pdf_path = tmp_path / "text.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "This PDF page contains enough visible text to avoid optical character recognition.")
    pdf.save(pdf_path)
    pdf.close()

    markdown_result = extract_file("md", "1", markdown, "text/markdown")
    html_result = extract_file("html", "1", html, "text/html")
    csv_result = extract_file("csv", "1", csv_path, "text/csv")
    json_result = extract_file("json", "1", json_path, "application/json")
    pdf_result = extract_file("pdf", "1", pdf_path, "application/pdf")

    assert [block.locator for block in markdown_result.blocks] == ["section:Setup", "section:Verify"]
    assert "SECRET" not in " ".join(block.text for block in html_result.blocks)
    assert csv_result.blocks[0].locator == "sheet:CSV!A1:B2"
    assert json_result.blocks[0].locator.startswith("json:$")
    assert pdf_result.blocks[0].locator == "page:1"


def test_html_without_heading_or_paragraph_uses_visible_text(tmp_path: Path):
    path = tmp_path / "simple.html"
    path.write_text("<div>Visible fallback body</div><script>secret()</script>", encoding="utf-8")

    result = extract_file("html", "1", path, "text/html")

    assert result.blocks[0].text == "Visible fallback body"


def test_xlsx_formats_date_cells(tmp_path: Path):
    path = tmp_path / "dates.xlsx"
    workbook = Workbook()
    workbook.active.append(["Due", date(2026, 7, 18)])
    workbook.save(path)

    result = extract_file("sheet", "1", path, XLSX_MIME)

    assert "2026-07-18" in result.blocks[0].text


def test_json_node_limit_is_enforced(tmp_path: Path):
    path = tmp_path / "large.json"
    path.write_text(json.dumps(list(range(20_001))), encoding="utf-8")

    with pytest.raises(DriveRagError, match="20,000"):
        extract_file("json", "1", path, "application/json")


def test_file_limit_and_unsupported_format_are_typed(tmp_path: Path):
    large = tmp_path / "large.txt"
    with large.open("wb") as stream:
        stream.truncate(250 * 1024 * 1024 + 1)
    unsupported = tmp_path / "payload.bin"
    unsupported.write_bytes(b"binary")

    with pytest.raises(DriveRagError, match="250 MiB"):
        extract_file("large", "1", large, "text/plain")
    with pytest.raises(DriveRagError, match="unsupported") as error:
        extract_file("bin", "1", unsupported, "application/octet-stream")
    assert error.value.code == "UNSUPPORTED_FORMAT"


@pytest.mark.parametrize("mime_type", [DOCX_MIME, XLSX_MIME, PPTX_MIME])
def test_corrupt_office_files_return_typed_extraction_error(tmp_path: Path, mime_type: str):
    path = tmp_path / "corrupt.office"
    path.write_bytes(b"not an office archive")

    with pytest.raises(DriveRagError) as error:
        extract_file("corrupt", "1", path, mime_type)

    assert error.value.code == "EXTRACTION_FAILED"


def _forbid_docx_parser(monkeypatch: pytest.MonkeyPatch):
    def forbidden_parser(*args, **kwargs):
        raise AssertionError("Office parser ran before archive preflight")

    monkeypatch.setattr(extract_module, "Document", forbidden_parser)


def test_office_archive_member_count_is_bounded_before_parser(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    path = tmp_path / "many.docx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_STORED) as archive:
        for index in range(3):
            archive.writestr(f"word/item-{index}.xml", b"x")
    monkeypatch.setattr(extract_module, "MAX_ARCHIVE_MEMBERS", 2, raising=False)
    _forbid_docx_parser(monkeypatch)

    def forbidden_zipfile(*args, **kwargs):
        raise AssertionError("full central directory was materialized")

    monkeypatch.setattr(extract_module, "ZipFile", forbidden_zipfile)

    with pytest.raises(DriveRagError, match="member count"):
        extract_file("doc", "1", path, DOCX_MIME)


def test_office_archive_does_not_trust_underreported_eocd_member_count(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    path = tmp_path / "underreported.docx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_STORED) as archive:
        for index in range(3):
            archive.writestr(f"word/item-{index}.xml", b"x")
    payload = bytearray(path.read_bytes())
    eocd_offset = payload.rfind(b"PK\x05\x06")
    assert eocd_offset >= 0
    struct.pack_into("<HH", payload, eocd_offset + 8, 1, 1)
    path.write_bytes(payload)
    monkeypatch.setattr(extract_module, "MAX_ARCHIVE_MEMBERS", 2, raising=False)

    def forbidden_zipfile(*args, **kwargs):
        raise AssertionError("full central directory was materialized")

    monkeypatch.setattr(extract_module, "ZipFile", forbidden_zipfile)

    with pytest.raises(DriveRagError, match="member count"):
        extract_file("doc", "1", path, DOCX_MIME)


def test_office_archive_central_directory_is_bounded_before_materialization(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    path = tmp_path / "directory.docx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_STORED) as archive:
        archive.writestr("word/document.xml", b"x")
    monkeypatch.setattr(
        extract_module, "MAX_ARCHIVE_CENTRAL_DIRECTORY_BYTES", 1, raising=False
    )

    def forbidden_zipfile(*args, **kwargs):
        raise AssertionError("full central directory was materialized")

    monkeypatch.setattr(extract_module, "ZipFile", forbidden_zipfile)

    with pytest.raises(DriveRagError, match="central directory"):
        extract_file("doc", "1", path, DOCX_MIME)


def test_office_archive_member_and_aggregate_expansion_are_bounded(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    per_member = tmp_path / "member.docx"
    with zipfile.ZipFile(per_member, "w", compression=zipfile.ZIP_STORED) as archive:
        archive.writestr("word/document.xml", b"SECRET_DOCUMENT_BODY")
    aggregate = tmp_path / "aggregate.docx"
    with zipfile.ZipFile(aggregate, "w", compression=zipfile.ZIP_STORED) as archive:
        archive.writestr("word/a.xml", b"1234567890")
        archive.writestr("word/b.xml", b"1234567890")
    _forbid_docx_parser(monkeypatch)

    monkeypatch.setattr(extract_module, "MAX_ARCHIVE_MEMBER_BYTES", 10, raising=False)
    with pytest.raises(DriveRagError, match="member expansion") as member_error:
        extract_file("doc", "1", per_member, DOCX_MIME)
    assert "SECRET_DOCUMENT_BODY" not in str(member_error.value)

    monkeypatch.setattr(extract_module, "MAX_ARCHIVE_MEMBER_BYTES", 100, raising=False)
    monkeypatch.setattr(extract_module, "MAX_ARCHIVE_UNCOMPRESSED_BYTES", 15, raising=False)
    with pytest.raises(DriveRagError, match="aggregate expansion"):
        extract_file("doc", "1", aggregate, DOCX_MIME)


def test_office_archive_compression_ratio_is_bounded_before_parser(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    path = tmp_path / "ratio.docx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("word/document.xml", b"A" * 10_000)
    monkeypatch.setattr(extract_module, "MAX_ARCHIVE_COMPRESSION_RATIO", 2, raising=False)
    monkeypatch.setattr(
        extract_module, "MIN_ARCHIVE_RATIO_CHECK_BYTES", 100, raising=False
    )
    _forbid_docx_parser(monkeypatch)

    with pytest.raises(DriveRagError, match="compression ratio"):
        extract_file("doc", "1", path, DOCX_MIME)


def test_small_highly_compressible_office_member_is_not_rejected(tmp_path: Path):
    path = tmp_path / "repetitive.docx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("word/document.xml", b"<p>policy</p>" * 70_000)

    extract_module._preflight_office_archive(path)


@pytest.mark.parametrize("unsafe_name", ["../escape.xml", "/absolute.xml", "word\\evil.xml"])
def test_office_archive_rejects_unsafe_member_paths(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch, unsafe_name: str
):
    path = tmp_path / "unsafe.docx"
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_STORED) as archive:
        archive.writestr(unsafe_name, b"x")
    _forbid_docx_parser(monkeypatch)

    with pytest.raises(DriveRagError, match="unsafe archive member path"):
        extract_file("doc", "1", path, DOCX_MIME)


def test_office_archive_rejects_symlink_members_before_parser(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
):
    path = tmp_path / "symlink.docx"
    member = zipfile.ZipInfo("word/link.xml")
    member.create_system = 3
    member.external_attr = (stat.S_IFLNK | 0o777) << 16
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(member, b"document.xml")
    _forbid_docx_parser(monkeypatch)

    with pytest.raises(DriveRagError, match="unsafe archive member type"):
        extract_file("doc", "1", path, DOCX_MIME)


def test_extract_rejects_symlink_payload(tmp_path: Path):
    target = tmp_path / "target.txt"
    target.write_text("private", encoding="utf-8")
    link = tmp_path / "link.txt"
    link.symlink_to(target)

    with pytest.raises(DriveRagError, match="symlink"):
        extract_file("file-a", "1", link, "text/plain")


def test_chunk_identity_includes_revision_locator_ordinal_and_text():
    first = ExtractedDocument(
        "file-a",
        "1",
        (ExtractedBlock("section:A", "alpha beta", {}),),
        "text",
    )
    revision = ExtractedDocument("file-a", "2", first.blocks, "text")
    locator = ExtractedDocument(
        "file-a",
        "1",
        (ExtractedBlock("section:B", "alpha beta", {}),),
        "text",
    )

    first_chunk = chunk_document(first, WordCounter(), 20, 4)[0]
    assert first_chunk.chunk_id != chunk_document(revision, WordCounter(), 20, 4)[0].chunk_id
    assert first_chunk.chunk_id != chunk_document(locator, WordCounter(), 20, 4)[0].chunk_id
    assert first_chunk.chunk_id.startswith("file-a:")
    assert len(first_chunk.chunk_id) == len("file-a:") + 64


def test_character_chunks_preserve_requested_overlap():
    document = ExtractedDocument(
        "file-a",
        "1",
        (ExtractedBlock("section:CJK", "abcdefghijklmnopqrstuv", {}),),
        "text",
    )

    chunks = chunk_document(document, CharacterCounter(), 10, 3)

    assert chunks[1].text.startswith(chunks[0].text[-3:])


def test_extract_cli_keeps_content_out_of_stdout_and_rejects_outside_paths(tmp_path: Path):
    state = tmp_path / "state"
    staging = state / "staging" / "run"
    staging.mkdir(parents=True)
    descriptor = staging / "descriptor.json"
    descriptor.write_text(
        json.dumps(
            {
                "schema_version": "1",
                "file_id": "file-a",
                "revision": "7",
                "mime_type": "text/plain",
                "native_kind": None,
            }
        ),
        encoding="utf-8",
    )
    payload = staging / "payload.txt"
    payload.write_text("TOP SECRET DOCUMENT BODY", encoding="utf-8")
    output = staging / "extracted.json"
    script = Path("/usr/local/share/codex-drive-rag/skills/drive-rag/scripts/drive_rag.py")

    completed = subprocess.run(
        [
            sys.executable,
            str(script),
            "--state-root",
            str(state),
            "extract",
            "--descriptor",
            str(descriptor),
            "--payload",
            str(payload),
            "--output",
            str(output),
        ],
        text=True,
        capture_output=True,
        check=False,
    )

    assert completed.returncode == 0, completed.stderr
    assert "TOP SECRET" not in completed.stdout
    result = json.loads(completed.stdout)
    assert result["schema_version"] == "1"
    assert result["counts"] == {"blocks": 1}
    assert read_json(output)["file_id"] == "file-a"

    outside = tmp_path / "outside.txt"
    outside.write_text("outside", encoding="utf-8")
    rejected = subprocess.run(
        [
            sys.executable,
            str(script),
            "--state-root",
            str(state),
            "extract",
            "--descriptor",
            str(descriptor),
            "--payload",
            str(outside),
            "--output",
            str(output),
        ],
        text=True,
        capture_output=True,
        check=False,
    )
    assert rejected.returncode == 2
    assert json.loads(rejected.stdout)["error"]["code"] == "UNSAFE_PATH"


def test_extract_cli_rejects_native_kind_mismatch(tmp_path: Path):
    state = tmp_path / "state"
    staging = state / "staging" / "run"
    staging.mkdir(parents=True)
    descriptor = staging / "descriptor.json"
    descriptor.write_text(
        json.dumps(
            {
                "schema_version": "1",
                "file_id": "file-a",
                "revision": "7",
                "mime_type": "application/vnd.google-apps.document",
                "native_kind": "document",
            }
        ),
        encoding="utf-8",
    )
    payload = staging / "payload.pdf"
    payload.write_bytes(b"%PDF-1.7\n")
    structured = staging / "structured.json"
    structured.write_text(
        json.dumps({"kind": "spreadsheet", "sheets": []}), encoding="utf-8"
    )
    output = staging / "output.json"
    script = Path("/usr/local/share/codex-drive-rag/skills/drive-rag/scripts/drive_rag.py")

    completed = subprocess.run(
        [
            sys.executable,
            str(script),
            "--state-root",
            str(state),
            "extract",
            "--descriptor",
            str(descriptor),
            "--payload",
            str(payload),
            "--structured",
            str(structured),
            "--output",
            str(output),
        ],
        text=True,
        capture_output=True,
        check=False,
    )

    assert completed.returncode == 2
    assert json.loads(completed.stdout)["error"]["code"] == "INVALID_DESCRIPTOR"
